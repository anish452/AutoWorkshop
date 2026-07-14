"""
Generate a professionally designed Progress Report 1 PowerPoint for AutoRepairAgent.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

BASE = Path(__file__).resolve().parent
SCREENSHOTS = BASE / "presentation-screenshots"
OUTPUT = BASE / "Progress_Report_1_AutoRepairAgent_v2.pptx"

# ── Professional colour palette ──────────────────────────────────────────────
NAVY       = RGBColor(15,  39,  68)    # #0F2744  primary dark
BLUE       = RGBColor(25, 118, 210)    # #1976D2  primary accent
SKY        = RGBColor(227, 242, 253)   # #E3F2FD  light accent
TEAL       = RGBColor(0,  150, 136)    # #009688  secondary accent
WHITE      = RGBColor(255, 255, 255)
OFF_WHITE  = RGBColor(248, 250, 252)   # #F8FAFC  slide bg
DARK       = RGBColor(33,  33,  33)    # #212121  body text
GREY       = RGBColor(97,  97,  97)    # #616161  muted text
LIGHT_GREY = RGBColor(189, 189, 189)  # borders
GREEN      = RGBColor(46, 125,  50)    # success
ORANGE     = RGBColor(230, 81,  0)      # warning
RED        = RGBColor(198, 40,  40)    # error

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
FOOTER_H = Inches(0.42)
HEADER_H = Inches(1.05)
MARGIN = Inches(0.55)

_slide_counter = 0


# ── Diagram generator ─────────────────────────────────────────────────────────
def create_architecture_diagram(path: Path) -> None:
    w, h = 1400, 820
    img = Image.new("RGB", (w, h), "#F8FAFC")
    draw = ImageDraw.Draw(img)

    # Header bar
    draw.rectangle([0, 0, w, 70], fill="#0F2744")
    draw.text((w // 2, 35), "AutoRepairAgent — Phase 1 System Architecture",
              fill="white", anchor="mm")

    def box(x1, y1, x2, y2, fill, border, title, sub=""):
        draw.rounded_rectangle([x1, y1, x2, y2], radius=14, fill=fill, outline=border, width=2)
        cx = (x1 + x2) // 2
        draw.text((cx, (y1 + y2) // 2 - 18), title, fill="#0F2744", anchor="mm")
        if sub:
            draw.text((cx, (y1 + y2) // 2 + 14), sub, fill="#616161", anchor="mm")

    def arrow_h(x1, y, x2):
        draw.line([x1, y, x2 - 14, y], fill="#1976D2", width=3)
        draw.polygon([(x2, y), (x2 - 14, y - 7), (x2 - 14, y + 7)], fill="#1976D2")

    def arrow_v(x, y1, y2):
        draw.line([x, y1, x, y2 - 14], fill="#1976D2", width=3)
        draw.polygon([(x, y2), (x - 7, y2 - 14), (x + 7, y2 - 14)], fill="#1976D2")

    # Top tier
    box(60,  110, 360,  230, "#E3F2FD", "#1976D2", "React Frontend", "Vite · MUI · React Router")
    box(530, 110, 870,  230, "#E8F5E9", "#388E3C", "Express.js API",  "JWT · RBAC · Zod")
    box(1040,110, 1340, 230, "#FFF3E0", "#F57C00", "PostgreSQL DB",   "Prisma ORM · Neon")

    arrow_h(360, 170, 530)
    arrow_h(870, 170, 1040)

    # Middle tier — clean architecture layers
    box(60,  310, 420,  420, "#F3E5F5", "#7B1FA2", "Presentation Layer",
        "Routes · Controllers · Middlewares")
    box(490, 310, 910,  420, "#E0F7FA", "#0097A7", "Application Layer",
        "AuthService · UserService · AuditService")
    box(980, 310, 1340, 420, "#ECEFF1", "#546E7A", "Infrastructure Layer",
        "Repositories · Prisma Client")

    arrow_v(700, 230, 280)
    arrow_v(200, 420, 490)
    arrow_v(700, 420, 490)
    arrow_v(1160, 420, 490)

    # Bottom — security row
    draw.rounded_rectangle([60, 490, 1340, 580], radius=14, fill="#E8EAF6", outline="#3F51B5", width=2)
    draw.text((700, 520), "Security Layer", fill="#1A237E", anchor="mm")
    draw.text((700, 552), "JWT Authentication  ·  bcrypt Hashing  ·  Helmet  ·  CORS  ·  Rate Limiting",
              fill="#616161", anchor="mm")

    # Footer citation
    draw.rectangle([0, h - 50, w, h], fill="#0F2744")
    draw.text((w // 2, h - 25),
              "Clean Architecture + Repository Pattern  (Martin, 2017)",
              fill="#90CAF9", anchor="mm")
    img.save(path)


# ── Slide helpers ─────────────────────────────────────────────────────────────
def _rect(slide, l, t, w, h, fill, line=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _text(slide, l, t, w, h, text, size=14, bold=False, color=DARK,
          align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def _apply_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OFF_WHITE


def _add_footer(slide, slide_num: int, section: str = ""):
    _rect(slide, Inches(0), SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, NAVY)
    _text(slide, MARGIN, SLIDE_H - FOOTER_H + Inches(0.07), Inches(5), Inches(0.3),
          "AutoRepairAgent  |  Progress Report 1", size=9, color=SKY)
    if section:
        _text(slide, Inches(3.5), SLIDE_H - FOOTER_H + Inches(0.07), Inches(3), Inches(0.3),
              section, size=9, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    _text(slide, Inches(8.5), SLIDE_H - FOOTER_H + Inches(0.07), Inches(1), Inches(0.3),
          str(slide_num), size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def _add_header(slide, title: str, subtitle: str = "", section_tag: str = ""):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, NAVY)
    # Accent stripe
    _rect(slide, Inches(0), HEADER_H, SLIDE_W, Inches(0.06), BLUE)
    # Logo circle
    logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.35), Inches(0.18), Inches(0.65), Inches(0.65))
    logo.fill.solid()
    logo.fill.fore_color.rgb = BLUE
    logo.line.fill.background()
    _text(slide, Inches(0.35), Inches(0.28), Inches(0.65), Inches(0.45),
          "AR", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Section tag pill
    if section_tag:
        tag = _rect(slide, Inches(8.1), Inches(0.22), Inches(1.55), Inches(0.38), BLUE)
        _text(slide, Inches(8.1), Inches(0.26), Inches(1.55), Inches(0.3),
              section_tag, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.15), Inches(0.15), Inches(7), Inches(0.55),
          title, size=24, bold=True, color=WHITE)
    if subtitle:
        _text(slide, Inches(1.15), Inches(0.65), Inches(7), Inches(0.35),
              subtitle, size=12, color=SKY)


def _content_card(slide, top=Inches(1.25), height=Inches(5.75)):
    card = _rect(slide, MARGIN, top, SLIDE_W - MARGIN * 2, height, WHITE, LIGHT_GREY, radius=True)
    card.shadow.inherit = False
    return card


def _bullets_in_card(slide, items, top=Inches(1.45), left=Inches(0.85),
                     width=Inches(8.5), height=Inches(5.4), size=15, bullet_color=BLUE):
    _content_card(slide, top=Inches(1.25), height=Inches(5.75))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Strip leading emoji markers for cleaner look; re-add styled prefix
        text = item
        if text.startswith("✅"):
            p.text = text
            p.font.color.rgb = GREEN
        elif text.startswith("❌") or text.startswith("⏳"):
            p.text = text
            p.font.color.rgb = ORANGE if text.startswith("⏳") else RED
        else:
            p.text = f"▸  {text}"
            p.font.color.rgb = DARK
        p.font.size = Pt(size)
        p.space_after = Pt(10)
        p.line_spacing = 1.3
    return box


def _new_slide(prs, title, subtitle="", section_tag="", section=""):
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_bg(slide)
    _add_header(slide, title, subtitle, section_tag)
    _add_footer(slide, _slide_counter, section)
    return slide


def _section_divider(prs, number: str, title: str, subtitle: str):
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Decorative circles
    for x, y, r, alpha in [(8.5, 1.0, 2.5, SKY), (1.0, 5.5, 1.8, BLUE), (9.0, 5.0, 1.2, TEAL)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
        c.fill.solid()
        c.fill.fore_color.rgb = alpha
        c.line.fill.background()

    _text(slide, MARGIN, Inches(2.5), Inches(2), Inches(0.8),
          number, size=60, bold=True, color=BLUE)
    _rect(slide, MARGIN, Inches(3.4), Inches(1.2), Inches(0.07), TEAL)
    _text(slide, MARGIN, Inches(3.6), Inches(8), Inches(1.0),
          title, size=36, bold=True, color=WHITE)
    _text(slide, MARGIN, Inches(4.6), Inches(8), Inches(0.6),
          subtitle, size=16, color=SKY)
    _add_footer(slide, _slide_counter)
    return slide


def _image_slide(prs, title, subtitle, image_path, caption, section_tag="", section=""):
    slide = _new_slide(prs, title, subtitle, section_tag, section)
    _content_card(slide, top=Inches(1.2), height=Inches(5.85))
    if image_path and image_path.exists():
        img = Image.open(image_path)
        iw, ih = img.size
        # Max display area inside the content card (inches)
        max_w_in = 8.6
        max_h_in = 5.2
        # Convert pixel dimensions to inches (96 DPI standard)
        img_w_in = iw / 96.0
        img_h_in = ih / 96.0
        scale = min(max_w_in / img_w_in, max_h_in / img_h_in)
        pic_w = Inches(img_w_in * scale)
        pic_h = Inches(img_h_in * scale)
        left = (SLIDE_W - pic_w) / 2
        top = Inches(1.35) + (Inches(5.2) - pic_h) / 2
        slide.shapes.add_picture(str(image_path), left, top, pic_w, pic_h)
    else:
        # Placeholder if image missing
        _text(slide, MARGIN, Inches(3.5), Inches(8.5), Inches(0.5),
              f"[Screenshot not found: {image_path}]", size=14, color=RED, align=PP_ALIGN.CENTER)
    if caption:
        _rect(slide, MARGIN, Inches(6.55), SLIDE_W - MARGIN * 2, Inches(0.42),
              RGBColor(241, 245, 249), LIGHT_GREY, radius=True)
        _text(slide, MARGIN + Inches(0.15), Inches(6.6), SLIDE_W - MARGIN * 2 - Inches(0.3),
              Inches(0.32), caption, size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    return slide


def _two_column_slide(prs, title, subtitle, left_items, right_items, left_title, right_title,
                      section_tag="", section=""):
    slide = _new_slide(prs, title, subtitle, section_tag, section)
    _content_card(slide)
    col_w = Inches(4.0)
    # Left column header
    _rect(slide, Inches(0.75), Inches(1.4), col_w, Inches(0.38), BLUE)
    _text(slide, Inches(0.75), Inches(1.44), col_w, Inches(0.32),
          left_title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Right column header
    _rect(slide, Inches(5.1), Inches(1.4), col_w, Inches(0.38), TEAL)
    _text(slide, Inches(5.1), Inches(1.44), col_w, Inches(0.32),
          right_title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for col_items, col_left, col_top in [
        (left_items, Inches(0.75), Inches(1.9)),
        (right_items, Inches(5.1), Inches(1.9)),
    ]:
        box = slide.shapes.add_textbox(col_left, col_top, col_w, Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(col_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.space_after = Pt(8)
    return slide


def _table_slide(prs, title, subtitle, headers, rows, col_widths, section_tag="", section=""):
    slide = _new_slide(prs, title, subtitle, section_tag, section)
    _content_card(slide)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(0.75), Inches(1.45),
                                       Inches(8.5), Inches(0.45 * n_rows))
    tbl = tbl_shape.table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)

    # Header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    status_colors = {"Complete": GREEN, "Pending": ORANGE, "Resolved": GREEN, "Open": RED}
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if ri % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = status_colors.get(val, DARK)
                p.alignment = PP_ALIGN.CENTER if ci == n_cols - 1 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    global _slide_counter
    _slide_counter = 0

    arch_path = SCREENSHOTS / "06-architecture-diagram.png"
    create_architecture_diagram(arch_path)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── SLIDE 1: Title ──────────────────────────────────────────────────────
    _slide_counter += 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Decorative elements
    _rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _rect(s, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), TEAL)
    for x, y, r in [(7.5, 0.5, 3.0), (-1.0, 4.5, 2.5), (8.5, 5.0, 1.5)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
        c.fill.solid()
        c.fill.fore_color.rgb = BLUE
        c.fill.transparency = 0.85
        c.line.fill.background()

    # Logo block
    logo_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(3.75), Inches(0.9), Inches(2.5), Inches(2.5))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = BLUE
    logo_bg.line.fill.background()
    _text(s, Inches(3.75), Inches(1.55), Inches(2.5), Inches(1.2),
          "AR", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _text(s, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.9),
          "AutoRepairAgent", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(4.45), Inches(8.4), Inches(0.55),
          "AI-Powered Vehicle Repair Workshop Management System",
          size=16, color=SKY, align=PP_ALIGN.CENTER)
    _rect(s, Inches(3.5), Inches(5.1), Inches(3), Inches(0.05), TEAL)
    _text(s, Inches(0.8), Inches(5.25), Inches(8.4), Inches(0.45),
          "Progress Report 1  —  Phase 1 Foundation Development",
          size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    info_lines = [
        ("Student", "[Your Name]  |  ID: [Student ID]"),
        ("Course",  "[Course Name / Code]"),
        ("College", "[College Name]"),
        ("Period",  "Reporting Period: [Start Date] – [End Date]"),
        ("Date",    "Submission Date: July 2026"),
    ]
    for i, (label, val) in enumerate(info_lines):
        y = Inches(5.85) + Inches(i * 0.28)
        _text(s, Inches(2.0), y, Inches(1.4), Inches(0.25), label + ":", size=10,
              bold=True, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)
        _text(s, Inches(3.5), y, Inches(5), Inches(0.25), val, size=10, color=WHITE)

    _add_footer(s, 1)

    # ── SLIDE 2: Agenda ───────────────────────────────────────────────────────
    _table_slide(prs,
        "Presentation Agenda",
        "Organisation of This Report — 20 Slides",
        ["Section", "Topics Covered", "Slides"],
        [
            ["Introduction",  "Project overview, objectives, technology stack",         "3 – 5"],
            ["Progress",      "Database, authentication, admin modules + screenshots",  "6 – 14"],
            ["Pending Work",  "Phase 2 modules not yet started",                        "15"],
            ["Challenges",    "Problems encountered & how they were resolved",          "16 – 17"],
            ["Status & Plan", "Project on-track assessment, next period plan",          "18"],
            ["References",    "APA 7th Edition bibliography",                         "19 – 20"],
        ],
        [1.8, 5.2, 1.5],
        section_tag="AGENDA",
    )

    # ── SLIDE 3: Introduction ─────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Project Introduction", "Background & Purpose of AutoRepairAgent",
                   "01 — INTRO", "Introduction"),
        [
            "AutoRepairAgent is a full-stack web system for vehicle repair workshop management.",
            "The system uses AI to classify customer complaints and route jobs to departments.",
            "Phase 1 focused on technical foundation: database, auth, admin, and UI design.",
            "Development followed structured software engineering practices (Sommerville, 2016).",
            "This report documents Phase 1 progress, challenges faced, and Phase 2 plans.",
        ],
    )

    # ── SLIDE 4: Objectives ───────────────────────────────────────────────────
    _table_slide(prs,
        "Phase 1 Objectives",
        "Planned Deliverables for the First Reporting Period",
        ["#", "Module / Feature", "Status"],
        [
            ["1", "Database design & PostgreSQL with Prisma ORM",          "Complete"],
            ["2", "JWT authentication & bcrypt password hashing",          "Complete"],
            ["3", "Role-Based Access Control (8 roles)",                   "Complete"],
            ["4", "Admin module — User creation & management",             "Complete"],
            ["5", "Department creation & management",                      "Complete"],
            ["6", "Front-end interface design (React + MUI)",              "Complete"],
            ["7", "Customer, Vehicle & AI modules",                      "Pending"],
        ],
        [0.6, 5.9, 1.5],
        section_tag="01 — INTRO",
        section="Introduction",
    )

    # ── SLIDE 5: Tech stack ───────────────────────────────────────────────────
    _two_column_slide(prs,
        "Technology Stack",
        "Tools & Frameworks — Phase 1 (Sommerville, 2016)",
        [
            "Node.js (LTS) runtime environment",
            "Express.js REST API framework (Express.js, n.d.)",
            "Prisma ORM + PostgreSQL database (Prisma Data, Inc., 2024)",
            "JWT + bcrypt authentication (OWASP Foundation, 2021)",
            "Zod input validation",
            "Winston logging",
            "Clean Architecture (Martin, 2017)",
        ],
        [
            "React 18 + Vite build tool",
            "Material UI v5 components (Meta Platforms, Inc., 2024)",
            "TanStack React Query",
            "React Router v6 navigation",
            "React Hook Form + Zod",
            "Axios HTTP client",
            "Recharts data visualisation",
        ],
        "Backend", "Frontend",
        section_tag="01 — INTRO",
        section="Introduction",
    )

    # ── SLIDE 6: Architecture diagram ─────────────────────────────────────────
    _image_slide(prs,
        "System Architecture",
        "Clean Architecture — 3-Layer Backend + React Frontend (Martin, 2017)",
        arch_path,
        "Figure 1. Phase 1 three-tier system architecture with security layer",
        section_tag="02 — PROGRESS",
        section="Progress",
    )

    # ── SLIDE 7: Database ───────────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Database Development",
                   "PostgreSQL Schema & Prisma ORM (Prisma Data, Inc., 2024)",
                   "02 — PROGRESS", "Progress"),
        [
            "Designed 10 relational tables: Role, Department, User, Customer, Vehicle, Job, etc.",
            "Created enums: RoleName, DepartmentCode, JobStatus, AuditAction",
            "Initial migration: prisma/migrations/20260613122721_init",
            "Seed script inserts 5 departments, 8 test users, and sample data",
            "Foreign key relationships pre-built for Phase 2 business modules",
            "Database hosted on PostgreSQL (Neon cloud) with SSL connection",
        ],
    )

    # ── SLIDE 8: Authentication ─────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Authentication Module",
                   "Secure Login & Role-Based Access Control (OWASP Foundation, 2021)",
                   "02 — PROGRESS", "Progress"),
        [
            "POST /api/auth/login — validates credentials, returns signed JWT token",
            "Passwords hashed with bcrypt (salt rounds) before database storage",
            "JWT middleware (auth.js) protects all private API routes",
            "RBAC middleware (rbac.js) enforces role-based endpoint access",
            "Frontend AuthContext stores token in localStorage (ara_token)",
            "Protected routes redirect unauthenticated users to /login",
            "Axios interceptor attaches token and handles 401 session expiry",
        ],
    )

    # ── SLIDE 9: Login screenshot ─────────────────────────────────────────────
    login_img = SCREENSHOTS / "01-login-page.png"
    if not login_img.exists():
        login_img = SCREENSHOTS / "page-2026-07-07T21-13-28-569Z.png"
    _image_slide(prs,
        "Authentication — Login Interface",
        "System Screenshot: Secure User Login Page",
        login_img,
        "Figure 2. AutoRepairAgent login page — JWT authentication with role-based access",
        section_tag="02 — PROGRESS",
        section="Progress",
    )

    # ── SLIDE 10: User management ───────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Admin Module — User Management",
                   "Create, Edit, Activate & Deactivate System Users",
                   "02 — PROGRESS", "Progress"),
        [
            "REST API: GET / POST / PUT / DELETE  /api/admin/users",
            "Admin assigns roles: ADMIN, JOB_ADVISOR, MECHANICAL, CUSTOMER, etc.",
            "Users linked to workshop departments where applicable",
            "Activate / deactivate accounts without permanent deletion",
            "Zod validation on all input fields — backend and frontend",
            "Audit logging records every user management action",
        ],
    )

    # ── SLIDE 11: Users screenshot ──────────────────────────────────────────────
    _image_slide(prs,
        "User Management Interface",
        "System Screenshot: Admin Users Page",
        SCREENSHOTS / "03-users-page.png",
        "Figure 3. User list showing name, email, role, department, and active status",
        section_tag="02 — PROGRESS",
        section="Progress",
    )

    # ── SLIDE 12: Departments ───────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Department Management",
                   "Workshop Department Configuration & CRUD Operations",
                   "02 — PROGRESS", "Progress"),
        [
            "REST API: GET / POST / PUT / DELETE  /api/admin/departments",
            "5 default departments seeded: Mechanical, Electrical, Body Repair, Paint, General Inspection",
            "Each department: name, unique code, description, active/inactive status",
            "Department staff linked via departmentId on User model",
            "Card-based UI layout with create, edit, and delete dialog forms",
        ],
    )

    # ── SLIDE 13: Departments screenshot ────────────────────────────────────────
    _image_slide(prs,
        "Department Management Interface",
        "System Screenshot: Departments Page",
        SCREENSHOTS / "04-departments-page.png",
        "Figure 4. Workshop department cards with descriptions and management actions",
        section_tag="02 — PROGRESS",
        section="Progress",
    )

    # ── SLIDE 14: Front-End Design ──────────────────────────────────────────────
    _two_column_slide(prs,
        "Front-End Design",
        "React UI Layout & Reusable Components (Meta Platforms, Inc., 2024)",
        [
            "MainLayout — sidebar nav, app bar, theme toggle",
            "Role-based menu items per user role",
            "PageHeader, StatusChip, KpiCard components",
            "Light / dark mode theme support",
            "Inter + Space Grotesk typography",
            "Blue gradient professional colour scheme",
        ],
        [
            "React Hook Form + Zod validation",
            "TanStack React Query data fetching",
            "Axios API service with interceptors",
            "React Router v6 protected routes",
            "React Toastify notifications",
            "Responsive MUI v5 grid layout",
        ],
        "Layout & Theme", "Data & Routing",
        section_tag="02 — PROGRESS",
        section="Progress",
    )

    # ── SLIDE 15: Dashboard screenshot ──────────────────────────────────────────
    dash_img = SCREENSHOTS / "05-admin-dashboard-full.png"
    if not dash_img.exists():
        dash_img = SCREENSHOTS / "02-admin-dashboard.png"
    _image_slide(prs,
        "Admin Dashboard",
        "System Screenshot: Dashboard UI with KPI Summary Cards",
        dash_img,
        "Figure 5. Admin dashboard — job statistics and system-wide overview layout",
        section_tag="02 — PROGRESS",
        section="Progress",
    )

    # ── SLIDE 16: Not completed ─────────────────────────────────────────────────
    _table_slide(prs,
        "Items Not Yet Completed",
        "Planned for Phase 2 — Next Reporting Period",
        ["Module", "Description", "Status"],
        [
            ["Customer Module",    "Register & manage workshop customers",           "Pending"],
            ["Vehicle Module",     "Register vehicles linked to customers",          "Pending"],
            ["AI Job Creation",    "Complaint submission & auto job generation",     "Pending"],
            ["AI Job Assignment",  "DeepSeek API integration & department routing",  "Pending"],
            ["AI Analysis Log",    "History of all AI classification decisions",     "Pending"],
            ["Deployment",         "Cloud hosting, Docker, CI/CD pipeline",          "Pending"],
        ],
        [2.2, 4.8, 1.5],
        section_tag="03 — PENDING",
        section="Pending Work",
    )

    # ── SLIDE 17: Problems & Solutions (combined) ───────────────────────────────
    _two_column_slide(prs,
        "Problems Encountered & Solutions",
        "Technical Challenges and How They Were Resolved (Sommerville, 2016)",
        [
            "Node modules — separate npm install for backend & frontend",
            "Database migration failures — incorrect DATABASE_URL",
            "Missing libraries — bcrypt, jsonwebtoken, zod errors",
            "CORS errors — React (5173) ↔ Express (3000) conflict",
            "Missing .env variables blocked server startup",
        ],
        [
            "✅ npm install in root + UI/autorepairagent folders",
            "✅ Corrected DATABASE_URL with sslmode=require",
            "✅ Installed all packages; updated package.json",
            "✅ Configured CORS + Vite proxy → localhost:3000",
            "✅ Created .env from .env.example; added dev:all script",
        ],
        "Problems Encountered", "Solutions Applied",
        section_tag="03 — CHALLENGES",
        section="Challenges",
    )

    # ── SLIDE 18: Status & next plan ──────────────────────────────────────────
    slide = _new_slide(prs,
        "Project Status & Next Period Plan",
        "Assessment: ON TRACK — Phase 2 Ready",
        "04 — STATUS", "Status & Plan")
    _content_card(slide)

    _rect(slide, Inches(0.85), Inches(1.45), Inches(2.2), Inches(0.55), GREEN, radius=True)
    _text(slide, Inches(0.85), Inches(1.5), Inches(2.2), Inches(0.45),
          "✔  ON TRACK", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    items = [
        "All Phase 1 objectives completed within the planned reporting period.",
        "Foundation ready: database, authentication, admin, departments, UI design.",
        "Phase 2 focus: Customer module → Vehicle module → AI integration.",
        "Phase 2 focus: AI job creation, assignment, and analysis log.",
        "Risk: DeepSeek API key availability — mitigated by keyword fallback classifier.",
    ]
    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.2), Inches(8.3), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(9)

    _rect(slide, Inches(0.85), Inches(5.5), Inches(8.3), Inches(1.35), SKY, BLUE, radius=True)
    _text(slide, Inches(1.0), Inches(5.58), Inches(8), Inches(0.3),
          "Phase 2 Timeline (Proposed)", size=12, bold=True, color=NAVY)
    phases = "Week 1: Customers  →  Week 2: Vehicles  →  Week 3: AI Integration  →  Week 4: Job Assignment & Testing"
    _text(slide, Inches(1.0), Inches(5.95), Inches(8), Inches(0.7), phases, size=12, color=DARK)

    # ── SLIDE 19: Screenshot summary ──────────────────────────────────────────
    _table_slide(prs,
        "System Screenshots Summary",
        "Phase 1 UI Evidence — Figures 2 to 5",
        ["Figure", "Screen", "Module Demonstrated"],
        [
            ["Figure 2", "Login Page",           "Authentication & JWT login"],
            ["Figure 3", "Users Page",           "Admin user management"],
            ["Figure 4", "Departments Page",     "Department creation & CRUD"],
            ["Figure 5", "Admin Dashboard",      "Front-end design & KPI layout"],
        ],
        [1.2, 2.5, 4.8],
        section_tag="02 — PROGRESS",
        section="Screenshots",
    )

    # ── SLIDE 20: References ──────────────────────────────────────────────────
    slide = _new_slide(prs, "References", "APA 7th Edition — Bibliography", "05 — REFERENCES")
    _content_card(slide)

    refs = [
        "Express.js. (n.d.). Express documentation. https://expressjs.com",
        "Martin, R. C. (2017). Clean architecture: A craftsman's guide to software structure and design. Prentice Hall.",
        "Meta Platforms, Inc. (2024). React documentation. https://react.dev",
        "OWASP Foundation. (2021). Authentication cheat sheet. https://cheatsheetseries.owasp.org/cheatsheets/Authentication_Cheat_Sheet.html",
        "Prisma Data, Inc. (2024). Prisma ORM documentation. https://www.prisma.io/docs",
        "Sommerville, I. (2016). Software engineering (10th ed.). Pearson.",
    ]
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(8.3), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}.    {ref}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.space_after = Pt(14)
        p.line_spacing = 1.4

    _text(slide, Inches(0.85), Inches(6.55), Inches(8.3), Inches(0.35),
          "Note: In-text citations throughout this presentation follow APA 7th edition author-date format.",
          size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    assert len(prs.slides) == 20, f"Expected 20 slides, got {len(prs.slides)}"

    # Verify embedded image sizes are visible (> 1 inch wide)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                w_in = shape.width / 914400
                print(f"  Slide {i}: image {w_in:.1f}\" x {shape.height / 914400:.1f}\"")

    prs.save(OUTPUT)
    print(f"Created : {OUTPUT}")
    print(f"Slides  : {len(prs.slides)}")
    print(f"Screenshots embedded:")
    for label, path in [
        ("Login", login_img),
        ("Users", SCREENSHOTS / "03-users-page.png"),
        ("Departments", SCREENSHOTS / "04-departments-page.png"),
        ("Dashboard", dash_img),
        ("Architecture", arch_path),
    ]:
        status = "OK" if path.exists() else "MISSING"
        print(f"  [{status}] {label}: {path.name}")


if __name__ == "__main__":
    main()
