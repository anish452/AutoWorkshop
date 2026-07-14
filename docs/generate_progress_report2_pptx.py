"""
Generate a professionally designed Progress Report 2 (Final) PowerPoint for AutoRepairAgent.
Same template as Report 1 — max 20 slides with UI screenshots.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from PIL import Image, ImageDraw

BASE = Path(__file__).resolve().parent
SCREENSHOTS = BASE / "presentation-screenshots-report2"
OUTPUT = BASE / "Progress_Report_2_AutoRepairAgent.pptx"
REPORT_LABEL = "Progress Report 2"

NAVY       = RGBColor(15,  39,  68)
BLUE       = RGBColor(25, 118, 210)
SKY        = RGBColor(227, 242, 253)
TEAL       = RGBColor(0,  150, 136)
WHITE      = RGBColor(255, 255, 255)
OFF_WHITE  = RGBColor(248, 250, 252)
DARK       = RGBColor(33,  33,  33)
GREY       = RGBColor(97,  97,  97)
LIGHT_GREY = RGBColor(189, 189, 189)
GREEN      = RGBColor(46, 125,  50)
ORANGE     = RGBColor(230, 81,  0)
RED        = RGBColor(198, 40,  40)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
FOOTER_H = Inches(0.42)
HEADER_H = Inches(1.05)
MARGIN = Inches(0.55)
_slide_counter = 0


def create_workflow_diagram(path: Path) -> None:
    w, h = 1400, 720
    img = Image.new("RGB", (w, h), "#F8FAFC")
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, w, 60], fill="#0F2744")
    draw.text((w // 2, 30), "AutoRepairAgent — End-to-End Workflow (Phase 2)",
              fill="white", anchor="mm")

    steps = [
        (80,  "Register\nCustomer",   "#E3F2FD", "#1976D2"),
        (280, "Register\nVehicle",    "#E8F5E9", "#388E3C"),
        (480, "Submit\nComplaint",    "#FFF3E0", "#F57C00"),
        (680, "AI Analysis\n(DeepSeek)", "#F3E5F5", "#7B1FA2"),
        (880, "Create\nJobs",         "#E0F7FA", "#0097A7"),
        (1080,"Auto\nAssign",         "#ECEFF1", "#546E7A"),
        (1280,"Start /\nComplete",    "#E8F5E9", "#2E7D32"),
    ]
    y1, y2 = 200, 310
    for i, (x, label, fill, border) in enumerate(steps):
        draw.rounded_rectangle([x, y1, x + 160, y2], radius=12, fill=fill, outline=border, width=2)
        draw.text((x + 80, 255), label, fill="#212121", anchor="mm")
        if i < len(steps) - 1:
            nx = steps[i + 1][0]
            draw.line([x + 160, 255, nx - 8, 255], fill="#1976D2", width=3)
            draw.polygon([(nx, 255), (nx - 10, 250), (nx - 10, 260)], fill="#1976D2")

    draw.rounded_rectangle([80, 400, 1320, 520], radius=14, fill="#E8EAF6", outline="#3F51B5", width=2)
    draw.text((700, 430), "AI Analysis Log  +  Role-Based Dashboards  +  Audit Trail",
              fill="#1A237E", anchor="mm")
    draw.text((700, 475),
              "Admin · Job Advisor · Department User · Customer dashboards",
              fill="#616161", anchor="mm")
    draw.rectangle([0, h - 45, w, h], fill="#0F2744")
    draw.text((w // 2, h - 22),
              "Figure 1. Complete repair workflow delivered in Phase 2",
              fill="#90CAF9", anchor="mm")
    img.save(path)


# ── Shared slide helpers (same template as Report 1) ─────────────────────────
def _rect(slide, l, t, w, h, fill, line=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if radius:
        shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _text(slide, l, t, w, h, text, size=14, bold=False, color=DARK,
          align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def _apply_slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE


def _add_footer(slide, slide_num, section=""):
    _rect(slide, Inches(0), SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, NAVY)
    _text(slide, MARGIN, SLIDE_H - FOOTER_H + Inches(0.07), Inches(5), Inches(0.3),
          f"AutoRepairAgent  |  {REPORT_LABEL}", size=9, color=SKY)
    if section:
        _text(slide, Inches(3.5), SLIDE_H - FOOTER_H + Inches(0.07), Inches(3), Inches(0.3),
              section, size=9, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    _text(slide, Inches(8.5), SLIDE_H - FOOTER_H + Inches(0.07), Inches(1), Inches(0.3),
          str(slide_num), size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def _add_header(slide, title, subtitle="", section_tag=""):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, NAVY)
    _rect(slide, Inches(0), HEADER_H, SLIDE_W, Inches(0.06), BLUE)
    logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.35), Inches(0.18), Inches(0.65), Inches(0.65))
    logo.fill.solid()
    logo.fill.fore_color.rgb = BLUE
    logo.line.fill.background()
    _text(slide, Inches(0.35), Inches(0.28), Inches(0.65), Inches(0.45),
          "AR", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if section_tag:
        _rect(slide, Inches(8.1), Inches(0.22), Inches(1.55), Inches(0.38), BLUE)
        _text(slide, Inches(8.1), Inches(0.26), Inches(1.55), Inches(0.3),
              section_tag, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.15), Inches(0.15), Inches(7), Inches(0.55), title, size=24, bold=True, color=WHITE)
    if subtitle:
        _text(slide, Inches(1.15), Inches(0.65), Inches(7), Inches(0.35), subtitle, size=12, color=SKY)


def _content_card(slide, top=Inches(1.25), height=Inches(5.75)):
    return _rect(slide, MARGIN, top, SLIDE_W - MARGIN * 2, height, WHITE, LIGHT_GREY, radius=True)


def _new_slide(prs, title, subtitle="", section_tag="", section=""):
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_bg(slide)
    _add_header(slide, title, subtitle, section_tag)
    _add_footer(slide, _slide_counter, section)
    return slide


def _bullets_in_card(slide, items, top=Inches(1.45), left=Inches(0.85),
                     width=Inches(8.5), height=Inches(5.4), size=15):
    _content_card(slide)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("✅"):
            p.text, p.font.color.rgb = item, GREEN
        elif item.startswith("❌") or item.startswith("⏳"):
            p.text, p.font.color.rgb = item, ORANGE if item.startswith("⏳") else RED
        else:
            p.text, p.font.color.rgb = f"▸  {item}", DARK
        p.font.size = Pt(size)
        p.space_after = Pt(10)


def _image_slide(prs, title, subtitle, image_path, caption, section_tag="", section=""):
    slide = _new_slide(prs, title, subtitle, section_tag, section)
    _content_card(slide, top=Inches(1.2), height=Inches(5.85))
    if image_path and image_path.exists():
        img = Image.open(image_path)
        iw, ih = img.size
        max_w_in, max_h_in = 8.6, 5.2
        img_w_in, img_h_in = iw / 96.0, ih / 96.0
        scale = min(max_w_in / img_w_in, max_h_in / img_h_in)
        pic_w, pic_h = Inches(img_w_in * scale), Inches(img_h_in * scale)
        left = (SLIDE_W - pic_w) / 2
        top = Inches(1.35) + (Inches(5.2) - pic_h) / 2
        slide.shapes.add_picture(str(image_path), left, top, pic_w, pic_h)
    else:
        _text(slide, MARGIN, Inches(3.5), Inches(8.5), Inches(0.5),
              f"[Screenshot not found: {image_path}]", size=14, color=RED, align=PP_ALIGN.CENTER)
    if caption:
        _rect(slide, MARGIN, Inches(6.55), SLIDE_W - MARGIN * 2, Inches(0.42),
              RGBColor(241, 245, 249), LIGHT_GREY, radius=True)
        _text(slide, MARGIN + Inches(0.15), Inches(6.6), SLIDE_W - MARGIN * 2 - Inches(0.3),
              Inches(0.32), caption, size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    return slide


def _two_column_slide(prs, title, subtitle, left_items, right_items, left_title, right_title,
                      section_tag="", section=""):
    slide = _new_slide(prs, title, subtitle, section_tag, section)
    _content_card(slide)
    col_w = Inches(4.0)
    _rect(slide, Inches(0.75), Inches(1.4), col_w, Inches(0.38), BLUE)
    _text(slide, Inches(0.75), Inches(1.44), col_w, Inches(0.32),
          left_title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(5.1), Inches(1.4), col_w, Inches(0.38), TEAL)
    _text(slide, Inches(5.1), Inches(1.44), col_w, Inches(0.32),
          right_title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for col_items, col_left in [(left_items, Inches(0.75)), (right_items, Inches(5.1))]:
        box = slide.shapes.add_textbox(col_left, Inches(1.9), col_w, Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(col_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = GREEN if item.startswith("✅") else DARK
            p.space_after = Pt(8)
    return slide


def _table_slide(prs, title, subtitle, headers, rows, col_widths, section_tag="", section=""):
    slide = _new_slide(prs, title, subtitle, section_tag, section)
    _content_card(slide)
    n_cols, n_rows = len(headers), len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.75), Inches(1.45),
                                       Inches(8.5), Inches(0.45 * n_rows))
    tbl = tbl_shape.table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)
    status_colors = {"Complete": GREEN, "Pending": ORANGE, "Resolved": GREEN, "Deferred": ORANGE}
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size, p.font.bold, p.font.color.rgb = Pt(11), True, WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if ri % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = status_colors.get(val, DARK)
                p.alignment = PP_ALIGN.CENTER if ci == len(row) - 1 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide


def main():
    global _slide_counter
    _slide_counter = 0
    SCREENSHOTS.mkdir(parents=True, exist_ok=True)

    workflow_path = SCREENSHOTS / "07-workflow-diagram.png"
    create_workflow_diagram(workflow_path)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── SLIDE 1: Title ──────────────────────────────────────────────────────
    _slide_counter += 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    _rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _rect(s, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), TEAL)
    logo_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(3.75), Inches(0.9), Inches(2.5), Inches(2.5))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = BLUE
    logo_bg.line.fill.background()
    _text(s, Inches(3.75), Inches(1.55), Inches(2.5), Inches(1.2),
          "AR", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.9),
          "AutoRepairAgent", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(4.45), Inches(8.4), Inches(0.55),
          "AI-Powered Vehicle Repair Workshop Management System",
          size=16, color=SKY, align=PP_ALIGN.CENTER)
    _rect(s, Inches(3.5), Inches(5.1), Inches(3), Inches(0.05), TEAL)
    _text(s, Inches(0.8), Inches(5.25), Inches(8.4), Inches(0.45),
          "Progress Report 2  —  Phase 2 Business Modules (Final Report)",
          size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    for i, (label, val) in enumerate([
        ("Student", "[Your Name]  |  ID: [Student ID]"),
        ("Course",  "[Course Name / Code]"),
        ("College", "[College Name]"),
        ("Period",  "Reporting Period: [Start Date] – [End Date]"),
        ("Date",    "Submission Date: July 2026"),
    ]):
        y = Inches(5.85) + Inches(i * 0.28)
        _text(s, Inches(2.0), y, Inches(1.4), Inches(0.25), label + ":", size=10,
              bold=True, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)
        _text(s, Inches(3.5), y, Inches(5), Inches(0.25), val, size=10, color=WHITE)
    _add_footer(s, 1)

    # ── SLIDE 2: Agenda ───────────────────────────────────────────────────────
    _table_slide(prs, "Presentation Agenda", "Organisation of This Report — 20 Slides",
        ["Section", "Topics Covered", "Slides"],
        [
            ["Recap",         "Phase 2 context & objectives",                    "3 – 4"],
            ["Progress",      "Customer, Vehicle, AI Jobs + UI screenshots",     "5 – 14"],
            ["Future Work",   "Deployment & enhancement recommendations",        "15"],
            ["Challenges",    "AI integration, DB & subscription problems",      "16"],
            ["Conclusion",    "Project complete — final assessment",             "17 – 18"],
            ["References",    "APA 7th Edition bibliography",                    "19 – 20"],
        ], [1.8, 5.2, 1.5], section_tag="AGENDA")

    # ── SLIDE 3: Introduction ─────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Project Introduction", "Phase 2 — Building on Report 1 Foundation",
                   "01 — INTRO", "Introduction"),
        [
            "Progress Report 1 delivered: database, auth, admin, departments, and UI design.",
            "Phase 2 focused on core business modules: customers, vehicles, and AI-powered jobs.",
            "DeepSeek AI integrated for intelligent complaint classification (OpenAI, 2024).",
            "Full end-to-end workflow now functional: complaint → AI → job → assign → complete.",
            "This is the final progress report — all development objectives are complete.",
        ])

    # ── SLIDE 4: Phase 2 Objectives ─────────────────────────────────────────────
    _table_slide(prs, "Phase 2 Objectives", "Deliverables Completed in Second Reporting Period",
        ["#", "Module / Feature", "Status"],
        [
            ["1", "Customer creation & management module",       "Complete"],
            ["2", "Vehicle registration module",                 "Complete"],
            ["3", "AI job creation & complaint analysis",        "Complete"],
            ["4", "AI integration — DeepSeek API",               "Complete"],
            ["5", "Automatic job assignment to departments",     "Complete"],
            ["6", "AI analysis log & role-based dashboards",       "Complete"],
            ["7", "Production deployment & multilingual UI",   "Deferred"],
        ], [0.6, 5.9, 1.5], section_tag="01 — INTRO", section="Introduction")

    # ── SLIDE 5: Customer Module ────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Customer Management Module",
                   "Register & Manage Workshop Customers (Prisma Data, Inc., 2024)",
                   "02 — PROGRESS", "Progress"),
        [
            "REST API: GET / POST / PUT / DELETE  /api/customers",
            "Fields: first name, last name, phone, email, address",
            "Accessible to ADMIN and JOB_ADVISOR roles via RBAC",
            "Zod validation on backend and React Hook Form on frontend",
            "Customers linked to vehicles and customer portal accounts",
            "Audit logging for all customer create, update, delete actions",
        ])

    # ── SLIDE 6: Customers Screenshot ───────────────────────────────────────────
    _image_slide(prs, "Customer Management Interface",
        "System Screenshot: Customers Page",
        SCREENSHOTS / "r2-customers.png",
        "Figure 2. Customer list with search, pagination, and CRUD actions",
        section_tag="02 — PROGRESS", section="Progress")

    # ── SLIDE 7: Vehicle Module ──────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Vehicle Registration Module",
                   "Register Vehicles Linked to Customers",
                   "02 — PROGRESS", "Progress"),
        [
            "REST API: GET / POST / PUT / DELETE  /api/vehicles",
            "Fields: registration number, chassis, make, model, year, customerId",
            "Unique registration number validation",
            "Customer dropdown selector on create/edit form",
            "Vehicles required before AI job creation (lookup by registration)",
            "Seed data: test vehicles ABC123 and XYZ789",
        ])

    # ── SLIDE 8: Vehicles Screenshot ────────────────────────────────────────────
    _image_slide(prs, "Vehicle Registration Interface",
        "System Screenshot: Vehicles Page",
        SCREENSHOTS / "r2-vehicles.png",
        "Figure 3. Vehicle list showing registration, make/model, and linked customer",
        section_tag="02 — PROGRESS", section="Progress")

    # ── SLIDE 9: AI Job Creation ─────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "AI Job Creation & DeepSeek Integration",
                   "Intelligent Complaint Classification (OpenAI, 2024)",
                   "02 — PROGRESS", "Progress"),
        [
            "POST /api/jobs/analyze — accepts vehicle registration + complaint text",
            "DeepSeekAIService calls DeepSeek Chat Completions API",
            "AI extracts individual issues and classifies into 5 departments",
            "Returns structured JSON: issue, department, confidence, explanation",
            "Keyword-based fallback classifier when API is unavailable",
            "One job created per identified issue with unique job number",
        ])

    # ── SLIDE 10: Create Job Screenshot ─────────────────────────────────────────
    _image_slide(prs, "AI Job Creation Interface",
        "System Screenshot: Create Job Page",
        SCREENSHOTS / "r2-create-job.png",
        "Figure 4. AI-powered complaint submission and job classification form",
        section_tag="02 — PROGRESS", section="Progress")

    # ── SLIDE 11: Job Assignment ─────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "Job Assignment & Lifecycle Management",
                   "Auto-Assign, Start, Complete Workflow",
                   "02 — PROGRESS", "Progress"),
        [
            "autoAssignJob() assigns job to first active user in target department",
            "Job status flow: PENDING → ASSIGNED → IN_PROGRESS → COMPLETED",
            "POST /api/jobs/:id/start and /api/jobs/:id/complete endpoints",
            "Department users see only jobs in their own department",
            "JobAdvisor and Admin can edit/delete jobs before work starts",
            "Time taken and completion comments recorded on job completion",
        ])

    # ── SLIDE 12: Jobs Screenshot ───────────────────────────────────────────────
    _image_slide(prs, "Job Management Interface",
        "System Screenshot: Jobs List Page",
        SCREENSHOTS / "r2-jobs.png",
        "Figure 5. Job list with status filters, department badges, and assignment",
        section_tag="02 — PROGRESS", section="Progress")

    # ── SLIDE 13: AI Analysis Log ────────────────────────────────────────────────
    _bullets_in_card(
        _new_slide(prs, "AI Analysis Log Module",
                   "Audit Trail of All AI Classification Decisions",
                   "02 — PROGRESS", "Progress"),
        [
            "AIAnalysisLog model stores complaint, raw AI response, and parsed issues",
            "Log entry created automatically on every POST /api/jobs/analyze call",
            "AIAnalysisPage displays all logs with confidence score progress bars",
            "Shows linked jobs created from each analysis session",
            "Accessible to ADMIN and JOB_ADVISOR for accuracy review",
            "Four role-specific dashboards completed: Admin, Advisor, Dept, Customer",
        ])

    # ── SLIDE 14: AI Analysis Screenshot ────────────────────────────────────────
    _image_slide(prs, "AI Analysis Log Interface",
        "System Screenshot: AI Analysis Logs Page",
        SCREENSHOTS / "r2-ai-analysis.png",
        "Figure 6. AI analysis history with confidence scores and created jobs",
        section_tag="02 — PROGRESS", section="Progress")

    # ── SLIDE 15: Future Work ───────────────────────────────────────────────────
    _table_slide(prs, "Future Recommendations",
        "Items Not Completed — Post-Project Enhancements",
        ["Enhancement", "Description", "Status"],
        [
            ["Deployment",           "Cloud hosting, Docker, CI/CD pipeline",          "Deferred"],
            ["Multilingual UI",        "Sinhala, Tamil & English language support",      "Deferred"],
            ["SMS / Email Alerts",     "Job status notifications to customers & staff",  "Deferred"],
            ["Automated Testing",      "Jest unit tests, CI/CD quality pipeline",        "Deferred"],
            ["Inventory Management",   "Spare parts tracking and basic invoicing",       "Deferred"],
        ], [2.2, 4.8, 1.5], section_tag="03 — FUTURE", section="Future Work")

    # ── SLIDE 16: Problems & Solutions ──────────────────────────────────────────
    _two_column_slide(prs, "Problems Encountered & Solutions",
        "Phase 2 Technical Challenges — All Resolved (Sommerville, 2016)",
        [
            "AI JSON parsing — markdown-wrapped responses broke JSON.parse()",
            "Invalid department codes from AI output",
            "DeepSeek API subscription quota exhausted (HTTP 402/429)",
            "Database migration sync issues during schema updates",
            "Cross-module integration bugs between new modules",
        ],
        [
            "✅ Added JSON cleanup + DEPARTMENT_MAP normalisation",
            "✅ Implemented keyword fallback classifier",
            "✅ Renewed DeepSeek subscription; updated API key",
            "✅ Reset migrations; updated seed with test data",
            "✅ Fixed Prisma relations and React Query cache invalidation",
        ],
        "Problems Encountered", "Solutions Applied",
        section_tag="03 — CHALLENGES", section="Challenges")

    # ── SLIDE 17: Project Complete ──────────────────────────────────────────────
    slide = _new_slide(prs, "Project Status — Development Complete",
        "Final Assessment: ALL OBJECTIVES ACHIEVED", "04 — STATUS", "Conclusion")
    _content_card(slide)
    _rect(slide, Inches(0.85), Inches(1.45), Inches(3.0), Inches(0.55), GREEN, radius=True)
    _text(slide, Inches(0.85), Inches(1.5), Inches(3.0), Inches(0.45),
          "✔  PROJECT COMPLETE", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        "All Phase 1 (Report 1) and Phase 2 (Report 2) modules delivered and tested.",
        "End-to-end workflow verified: customer → vehicle → complaint → AI → job → complete.",
        "System demonstrable locally with seed accounts (admin@autorepair.com / Password123!).",
        "Future work (deployment, multilingual, notifications) documented as recommendations.",
        "This is the final progress report for the AutoRepairAgent college project.",
    ]
    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.2), Inches(8.3), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(9)
    _rect(slide, Inches(0.85), Inches(5.5), Inches(8.3), Inches(1.35), SKY, BLUE, radius=True)
    _text(slide, Inches(1.0), Inches(5.58), Inches(8), Inches(0.3),
          "Report 1 vs Report 2 Summary", size=12, bold=True, color=NAVY)
    _text(slide, Inches(1.0), Inches(5.95), Inches(8), Inches(0.7),
          "Report 1: Foundation (DB, Auth, Admin, UI)  →  "
          "Report 2: Business Logic (Customer, Vehicle, AI Jobs, Assignment, Logs)",
          size=12, color=DARK)

    # ── SLIDE 18: Workflow diagram ──────────────────────────────────────────────
    _image_slide(prs, "End-to-End System Workflow",
        "Complete Repair Process — Phase 2 Delivered",
        workflow_path,
        "Figure 1. Full workflow from customer registration to job completion",
        section_tag="04 — WORKFLOW", section="Workflow")

    # ── SLIDE 19: Screenshot summary ────────────────────────────────────────────
    _table_slide(prs, "System Screenshots Summary",
        "Phase 2 UI Evidence — Figures 2 to 6",
        ["Figure", "Screen", "Module Demonstrated"],
        [
            ["Figure 2", "Customers Page",    "Customer creation & management"],
            ["Figure 3", "Vehicles Page",     "Vehicle registration module"],
            ["Figure 4", "Create Job Page",     "AI complaint analysis & job creation"],
            ["Figure 5", "Jobs List Page",      "Job assignment & lifecycle management"],
            ["Figure 6", "AI Analysis Logs",    "AI decision audit trail"],
        ], [1.2, 2.5, 4.8], section_tag="02 — PROGRESS", section="Screenshots")

    # ── SLIDE 20: References ────────────────────────────────────────────────────
    slide = _new_slide(prs, "References", "APA 7th Edition — Bibliography", "05 — REFERENCES")
    _content_card(slide)
    refs = [
        "Express.js. (n.d.). Express documentation. https://expressjs.com",
        "Martin, R. C. (2017). Clean architecture: A craftsman's guide to software structure and design. Prentice Hall.",
        "Meta Platforms, Inc. (2024). React documentation. https://react.dev",
        "OpenAI. (2024). DeepSeek API documentation. https://api-docs.deepseek.com",
        "OWASP Foundation. (2021). Authentication cheat sheet. https://cheatsheetseries.owasp.org/cheatsheets/Authentication_Cheat_Sheet.html",
        "Prisma Data, Inc. (2024). Prisma ORM documentation. https://www.prisma.io/docs",
        "Sommerville, I. (2016). Software engineering (10th ed.). Pearson.",
    ]
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(8.3), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}.    {ref}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)
    _text(slide, Inches(0.85), Inches(6.55), Inches(8.3), Inches(0.35),
          "Note: In-text citations follow APA 7th edition author-date format.",
          size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    assert len(prs.slides) == 20, f"Expected 20 slides, got {len(prs.slides)}"
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"  Slide {i}: image {shape.width / 914400:.1f}\" wide")

    prs.save(OUTPUT)
    print(f"Created : {OUTPUT}")
    print(f"Slides  : {len(prs.slides)}")
    for label, name in [
        ("Customers", "r2-customers.png"), ("Vehicles", "r2-vehicles.png"),
        ("Create Job", "r2-create-job.png"), ("Jobs", "r2-jobs.png"),
        ("AI Analysis", "r2-ai-analysis.png"), ("Workflow", "07-workflow-diagram.png"),
    ]:
        p = SCREENSHOTS / name
        print(f"  [{'OK' if p.exists() else 'MISSING'}] {label}: {name}")


if __name__ == "__main__":
    main()
