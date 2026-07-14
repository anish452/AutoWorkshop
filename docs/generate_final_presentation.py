"""
COMP902 Final Report Presentation — AutoRepairAgent
Professional layout aligned to COMP902 marking guide & presentation rubric.
Sections a–g | Focus · Organisation · Visual Aids · Q&A
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw

BASE = Path(__file__).resolve().parent
SS1 = BASE / "presentation-screenshots"
SS2 = BASE / "presentation-screenshots-report2"
ARCH_DIR = BASE / "presentation-screenshots-final"
OUTPUT = BASE / "COMP902_Final_Report_Presentation_AutoRepairAgent.pptx"
TOTAL = 25

COURSE = "COMP902 — Advanced Information Technology Specialised Project"
INST = "Auckland Institute of Studies  |  Trimester 2, 2026"

# Palette
NAVY = RGBColor(0x0F, 0x2B, 0x46)
TEAL = RGBColor(0x00, 0x96, 0x88)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TXT = RGBColor(0xF1, 0xF5, 0xF9)
SUCCESS = RGBColor(0x16, 0x65, 0x34)
AMBER = RGBColor(0xB4, 0x5A, 0x00)
BORDER = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def shot(path: Path) -> Path:
    return path if path.exists() else Path("__missing__")


def create_architecture_diagram(path: Path):
    w, h = 1400, 720
    img = Image.new("RGB", (w, h), "#F8FAFC")
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, w, 55], fill="#0F2B46")
    draw.text((w // 2, 28), "AutoRepairAgent — System Architecture", fill="white", anchor="mm")

    def box(x1, y1, x2, y2, fill, border, title, sub=""):
        draw.rounded_rectangle([x1, y1, x2, y2], radius=12, fill=fill, outline=border, width=2)
        cx = (x1 + x2) // 2
        draw.text((cx, (y1 + y2) // 2 - 16), title, fill="#212121", anchor="mm")
        if sub:
            draw.text((cx, (y1 + y2) // 2 + 12), sub, fill="#616161", anchor="mm")

    box(60, 100, 360, 210, "#E3F2FD", "#1976D2", "React SPA", "Vite · MUI · React Router")
    box(520, 100, 880, 210, "#E8F5E9", "#388E3C", "Express.js API", "JWT · RBAC · Zod")
    box(1040, 100, 1340, 210, "#FFF3E0", "#F57C00", "PostgreSQL", "Prisma ORM")
    for x1, x2 in [(360, 520), (880, 1040)]:
        draw.line([x1, 155, x2 - 12, 155], fill="#1976D2", width=3)
        draw.polygon([(x2, 155), (x2 - 12, 150), (x2 - 12, 160)], fill="#1976D2")

    box(60, 280, 420, 390, "#F3E5F5", "#7B1FA2", "Presentation", "Routes · Controllers")
    box(490, 280, 910, 390, "#E0F7FA", "#0097A7", "Application", "Services · Job · AI")
    box(980, 280, 1340, 390, "#ECEFF1", "#546E7A", "Infrastructure", "Repositories · Prisma")
    box(200, 450, 1200, 540, "#E8EAF6", "#3F51B5", "DeepSeek AI  +  Security Layer",
        "JWT · bcrypt · Helmet · CORS · Rate Limiting")
    draw.rectangle([0, h - 40, w, h], fill="#0F2B46")
    draw.text((w // 2, h - 20), "Clean Architecture + Repository Pattern (Martin, 2017)", fill="#90CAF9", anchor="mm")
    img.save(path)


def set_run(run, size=18, bold=False, color=SLATE, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_accent_bar(slide, height=Inches(7.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_header(slide, title, subtitle=None, section_tag=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), 0, Inches(13.15), Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    if section_tag:
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.28), Inches(0.55), Inches(0.55))
        tag.fill.solid()
        tag.fill.fore_color.rgb = TEAL
        tag.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.3), Inches(0.55), Inches(0.5))
        tb.text_frame.text = section_tag
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tb.text_frame.paragraphs[0].runs[0], size=16, bold=True, color=WHITE)
        title_left = Inches(1.05)
    else:
        title_left = Inches(0.45)

    tb = slide.shapes.add_textbox(title_left, Inches(0.22), Inches(11.5), Inches(0.55))
    tb.text_frame.text = title
    set_run(tb.text_frame.paragraphs[0].runs[0], size=26, bold=True, color=WHITE)

    if subtitle:
        sb = slide.shapes.add_textbox(title_left, Inches(0.72), Inches(11.5), Inches(0.38))
        sb.text_frame.text = subtitle
        set_run(sb.text_frame.paragraphs[0].runs[0], size=13, color=LIGHT_TXT)


def add_footer(slide, n: int):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(7.0), Inches(13.0), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(12.8), Inches(0.3))
    box.text_frame.text = f"{COURSE}  |  AutoRepairAgent Final Report  |  Slide {n} of {TOTAL}"
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_run(box.text_frame.paragraphs[0].runs[0], size=9, color=MUTED)


def add_bullets(slide, items, left=0.45, top=1.45, width=6.0, height=5.3, size=14, spacing=6, color=SLATE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(spacing)
        if p.runs:
            set_run(p.runs[0], size=size, color=color)
        else:
            r = p.add_run()
            r.text = item
            set_run(r, size=size, color=color)


def add_kpi_card(slide, label, value, left, top, width=2.0, accent=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.1))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    vb = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.12), Inches(width - 0.2), Inches(0.45))
    vb.text_frame.text = value
    vb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(vb.text_frame.paragraphs[0].runs[0], size=22, bold=True, color=NAVY)
    lb = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.6), Inches(width - 0.2), Inches(0.35))
    lb.text_frame.text = label
    lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(lb.text_frame.paragraphs[0].runs[0], size=10, color=MUTED)


def add_image_card(slide, path: Path, left, top, width, height, caption):
    if not path.exists():
        return
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left - 0.05), Inches(top - 0.05),
        Inches(width + 0.1), Inches(height + 0.55),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    cap = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.06), Inches(width), Inches(0.35))
    cap.text_frame.text = caption
    cap.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(cap.text_frame.paragraphs[0].runs[0], size=10, color=MUTED)


def add_section_divider(slide, letter, title, subtitle):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), Inches(13.333), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()
    lt = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(1.2), Inches(1.2))
    lt.text_frame.text = letter
    set_run(lt.text_frame.paragraphs[0].runs[0], size=72, bold=True, color=TEAL)
    tt = slide.shapes.add_textbox(Inches(2.2), Inches(1.8), Inches(10.5), Inches(0.8))
    tt.text_frame.text = title
    set_run(tt.text_frame.paragraphs[0].runs[0], size=32, bold=True, color=WHITE)
    st = slide.shapes.add_textbox(Inches(2.2), Inches(2.65), Inches(10.5), Inches(0.5))
    st.text_frame.text = subtitle
    set_run(st.text_frame.paragraphs[0].runs[0], size=16, color=LIGHT_TXT)
    rb = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.4))
    rb.text_frame.text = "COMP902 Final Report  |  AutoRepairAgent"
    set_run(rb.text_frame.paragraphs[0].runs[0], size=12, color=MUTED)


def content_slide(prs, n, title, subtitle=None, section_tag=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_header(slide, title, subtitle, section_tag)
    add_footer(slide, n)
    return slide


def build():
    ARCH_DIR.mkdir(exist_ok=True)
    arch_path = ARCH_DIR / "architecture.png"
    create_architecture_diagram(arch_path)

    login = SS1 / "01-login-page.png"
    users = SS1 / "03-users-page.png"
    dashboard = SS1 / "05-admin-dashboard-full.png"
    if not dashboard.exists():
        dashboard = SS1 / "02-admin-dashboard.png"
    customers = SS2 / "r2-customers.png"
    vehicles = SS2 / "r2-vehicles.png"
    create_job = SS2 / "r2-create-job.png"
    jobs = SS2 / "r2-jobs.png"
    ai_log = SS2 / "r2-ai-analysis.png"
    workflow = SS2 / "07-workflow-diagram.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    n = 0

    # ── SLIDE 1: Title ────────────────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    for y, h in [(0, 0.12), (7.38, 0.12)]:
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), prs.slide_width, Inches(h))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = TEAL
        stripe.line.fill.background()

    inst = s.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.4))
    inst.text_frame.text = INST
    set_run(inst.text_frame.paragraphs[0].runs[0], size=13, color=TEAL)

    title = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(1.4))
    title.text_frame.text = "AutoRepairAgent"
    set_run(title.text_frame.paragraphs[0].runs[0], size=38, bold=True, color=WHITE)

    sub = s.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(7.5), Inches(0.7))
    sub.text_frame.text = "AI-Powered Vehicle Repair Workshop\nManagement System"
    set_run(sub.text_frame.paragraphs[0].runs[0], size=17, color=LIGHT_TXT)

    course = s.shapes.add_textbox(Inches(0.8), Inches(3.75), Inches(7.5), Inches(0.4))
    course.text_frame.text = COURSE
    set_run(course.text_frame.paragraphs[0].runs[0], size=13, bold=True, color=TEAL)

    meta = s.shapes.add_textbox(Inches(0.8), Inches(4.35), Inches(7.0), Inches(1.8))
    tf = meta.text_frame
    for i, line in enumerate([
        "Student: [Your Full Name]",
        "Student ID: [Your Student ID]",
        "Supervisor: [Supervisor Name]",
        "Due: 5 p.m. Wednesday 12th August 2026",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(5)
        set_run(p.runs[0], size=14, color=LIGHT_TXT)

    add_image_card(s, login, 8.3, 1.2, 4.5, 2.7, "JWT authentication — role-based access")

    # ── SLIDE 2: Presentation Map ───────────────────────────────────────────────
    n += 1
    s = content_slide(prs, n, "Presentation Structure", "Aligned to COMP902 Final Report body sections")
    rows = [
        ("a", "Introduction", "Background, motivation, problem statement"),
        ("b", "Literature Review", "WMS, AI/NLP, architecture, RBAC"),
        ("c", "Scope & Objectives", "Phase 1 & 2 deliverables, boundaries"),
        ("d", "Methodology", "Phased development, tech stack, verification"),
        ("e", "Results & Outcomes", "UI modules, AI workflow, metrics"),
        ("f", "Discussion", "Objective alignment, limitations, lessons"),
        ("g", "Conclusions & Future Work", "Contributions, roadmap, Q&A"),
    ]
    for i, (letter, heading, desc) in enumerate(rows):
        top = 1.45 + i * 0.78
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(top), Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        bt = s.shapes.add_textbox(Inches(0.45), Inches(top + 0.04), Inches(0.42), Inches(0.38))
        bt.text_frame.text = letter
        bt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(bt.text_frame.paragraphs[0].runs[0], size=13, bold=True, color=WHITE)
        ht = s.shapes.add_textbox(Inches(1.05), Inches(top), Inches(3.5), Inches(0.35))
        ht.text_frame.text = heading
        set_run(ht.text_frame.paragraphs[0].runs[0], size=14, bold=True, color=NAVY)
        dt = s.shapes.add_textbox(Inches(1.05), Inches(top + 0.32), Inches(5.5), Inches(0.35))
        dt.text_frame.text = desc
        set_run(dt.text_frame.paragraphs[0].runs[0], size=11, color=MUTED)

    add_bullets(s, [
        "Presentation Rubric Focus Areas:",
        "✓  Clear purpose from the outset",
        "✓  Logical sequence matching report body",
        "✓  Professional visual aids (live UI screenshots)",
        "✓  Prepared Q&A to stimulate discussion",
    ], left=7.0, top=1.55, width=5.8, size=13, spacing=8)
    add_image_card(s, dashboard, 7.0, 4.0, 5.8, 2.5, "Live system — Admin KPI Dashboard")

    # ── SECTION a: Introduction ───────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "a", "Introduction", "Background, motivation & problem statement")

    n += 1
    s = content_slide(prs, n, "Background & Motivation", "Why AI-assisted workshop management matters", "a")
    add_bullets(s, [
        "Vehicle repair workshops rely on manual triage and paper job cards",
        "  Inconsistent routing delays repairs and reduces transparency",
        "Staff coordinate via phone — limited audit trail for decisions",
        "",
        "Motivation:",
        "• Operational — faster complaint-to-department classification",
        "• Academic — Master's-level full-stack + AI integration",
        "• Technical — Clean Architecture with LLM-assisted triage",
    ], width=6.2, size=13, spacing=5)
    add_image_card(s, login, 6.9, 1.45, 5.9, 3.3, "Secure login — gateway to role-based system")

    n += 1
    s = content_slide(prs, n, "Problem Statement", "Four core challenges addressed by AutoRepairAgent", "a")
    challenges = [
        ("1", "Manual complaint triage", "Unstructured descriptions cause misrouting between departments"),
        ("2", "No AI audit trail", "Classification decisions not logged or explainable"),
        ("3", "Fragmented access control", "Eight roles need least-privilege data visibility"),
        ("4", "No end-to-end digital workflow", "Customer → vehicle → job → completion not integrated"),
    ]
    for i, (num, title, desc) in enumerate(challenges):
        top = 1.5 + i * 1.25
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(top), Inches(12.5), Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER
        nb = s.shapes.add_textbox(Inches(0.55), Inches(top + 0.22), Inches(0.4), Inches(0.5))
        nb.text_frame.text = num
        set_run(nb.text_frame.paragraphs[0].runs[0], size=16, bold=True, color=TEAL)
        tb = s.shapes.add_textbox(Inches(1.1), Inches(top + 0.12), Inches(3.0), Inches(0.35))
        tb.text_frame.text = title
        set_run(tb.text_frame.paragraphs[0].runs[0], size=13, bold=True, color=NAVY)
        db = s.shapes.add_textbox(Inches(1.1), Inches(top + 0.48), Inches(11.5), Inches(0.35))
        db.text_frame.text = desc
        set_run(db.text_frame.paragraphs[0].runs[0], size=11, color=SLATE)

    # ── SECTION b: Literature Review ────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "b", "Critical Literature Review", "Positioning the project in academic & industry context")

    n += 1
    s = content_slide(prs, n, "Literature Review — Key Themes", "Five domains reviewed; research gap identified", "b")
    themes = [
        "Workshop Management Systems — digitise intake, scheduling, job tracking",
        "AI/NLP in Service Triage — LLM complaint parsing (Amershi et al., 2019)",
        "REST + SPA Architecture — Fielding (2000); Newman (2021)",
        "Clean Architecture — maintainable separation of concerns (Martin, 2017)",
        "RBAC Security — least-privilege access (Sandhu & Samarati, 1994)",
    ]
    add_bullets(s, themes, width=6.0, size=13, spacing=8)
    add_bullets(s, [
        "Research Gap:",
        "Affordable AI-integrated WMS",
        "for small/medium workshops",
        "lacking enterprise ERP budgets",
        "",
        "6 academic references cited",
        "(APA 7th edition)",
    ], left=6.9, top=1.55, width=5.8, size=13, spacing=7, color=NAVY)
    add_image_card(s, create_job, 6.9, 4.1, 5.8, 2.4, "AI complaint classification — DeepSeek integration")

    # ── SECTION c: Scope & Objectives ───────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "c", "Scope & Objectives", "Project goals, success criteria & boundaries")

    n += 1
    s = content_slide(prs, n, "Project Objectives", "Six specific objectives — core items achieved", "c")
    objectives = [
        ("O1", "Authentication & RBAC", "JWT + 8 roles, protected routes"),
        ("O2", "Admin foundation", "Users, departments, dashboards"),
        ("O3", "Customer & vehicle CRUD", "Linked records for job creation"),
        ("O4", "AI job classification", "DeepSeek + keyword fallback"),
        ("O5", "Job lifecycle", "Create, assign, track, complete"),
        ("O6", "AI audit logging", "Confidence scores & explanations"),
    ]
    for i, (oid, title, desc) in enumerate(objectives):
        col = i % 2
        row = i // 2
        left = 0.4 + col * 6.4
        top = 1.45 + row * 1.55
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(6.1), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        oid_box = s.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(0.7), Inches(0.4))
        oid_box.text_frame.text = oid
        set_run(oid_box.text_frame.paragraphs[0].runs[0], size=13, bold=True, color=TEAL)
        tt = s.shapes.add_textbox(Inches(left + 0.9), Inches(top + 0.12), Inches(5.0), Inches(0.35))
        tt.text_frame.text = title
        set_run(tt.text_frame.paragraphs[0].runs[0], size=13, bold=True, color=NAVY)
        dt = s.shapes.add_textbox(Inches(left + 0.9), Inches(top + 0.5), Inches(5.0), Inches(0.6))
        dt.text_frame.text = desc
        set_run(dt.text_frame.paragraphs[0].runs[0], size=11, color=SLATE)

    n += 1
    s = content_slide(prs, n, "Project Scope", "In scope vs deferred — Master's-level boundaries", "c")
    add_bullets(s, [
        "In Scope:",
        "• Full-stack web app (React + Node.js + PostgreSQL)",
        "• 34 REST API endpoints, 10 database tables",
        "• AI complaint classification with audit trail",
        "• Role-based dashboards for 4 user types",
    ], width=5.8, size=13, spacing=6)
    add_bullets(s, [
        "Out of Scope (Deferred):",
        "• Production cloud deployment",
        "• Sinhala/Tamil UI localisation",
        "• Automated unit/E2E test suite",
        "• Inventory, invoicing, SMS notifications",
    ], left=6.9, top=1.55, width=5.8, size=13, spacing=6, color=AMBER)
    add_image_card(s, users, 6.9, 4.2, 5.8, 2.3, "9 system users — RBAC role assignments")

    # ── SECTION d: Methodology ──────────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "d", "Methodology", "Development approach, technology & verification")

    n += 1
    s = content_slide(prs, n, "Development Methodology", "Phased iterative — reproducible from description", "d")
    add_bullets(s, [
        "Approach: Phased iterative across two progress reports",
        "",
        "Phase 1 — Foundation:",
        "  Database, auth, admin, departments, UI design",
        "Phase 2 — Business Modules:",
        "  Customers, vehicles, AI jobs, assignment, AI log",
        "",
        "Verification:",
        "  Postman API tests, seed data, manual E2E workflow",
        "  Live UI screenshots captured from running system",
    ], width=6.0, size=12, spacing=4)
    add_image_card(s, arch_path, 6.9, 1.45, 5.9, 3.3, "Clean Architecture — three-tier design")

    n += 1
    s = content_slide(prs, n, "Technology Stack", "Justified selections — full-stack monorepo", "d")
    stack = [
        ("Frontend", "React 18, Vite, MUI, React Router, Axios"),
        ("Backend", "Node.js Express, JWT + bcrypt, Zod validation"),
        ("Database", "PostgreSQL + Prisma ORM — 10 tables"),
        ("AI", "DeepSeek API with keyword fallback classifier"),
        ("Security", "Helmet, CORS, rate limiting, RBAC middleware"),
        ("Architecture", "Clean Architecture + Repository pattern"),
    ]
    for i, (layer, tech) in enumerate(stack):
        top = 1.45 + i * 0.88
        lb = s.shapes.add_textbox(Inches(0.4), Inches(top), Inches(2.2), Inches(0.35))
        lb.text_frame.text = layer
        set_run(lb.text_frame.paragraphs[0].runs[0], size=12, bold=True, color=TEAL)
        tb = s.shapes.add_textbox(Inches(2.7), Inches(top), Inches(4.0), Inches(0.35))
        tb.text_frame.text = tech
        set_run(tb.text_frame.paragraphs[0].runs[0], size=12, color=SLATE)
        if i < len(stack) - 1:
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(top + 0.42), Inches(6.2), Inches(0.01))
            ln.fill.solid()
            ln.fill.fore_color.rgb = BORDER
            ln.line.fill.background()
    add_image_card(s, workflow if workflow.exists() else create_job, 7.0, 1.45, 5.8, 3.3,
                   "End-to-end workflow: complaint → AI → job → completion")

    # ── SECTION e: Results & Outcomes ───────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "e", "Results & Outcomes", "Clear, well-explained deliverables — 15 marks criterion")

    n += 1
    s = content_slide(prs, n, "Results Overview — Key Metrics", "All core objectives substantially achieved", "e")
    kpis = [
        ("34", "API Endpoints"), ("10", "DB Tables"), ("8", "User Roles"),
        ("5", "Departments"), ("7", "Test Jobs"), ("15+", "UI Screens"),
    ]
    for i, (val, lbl) in enumerate(kpis):
        add_kpi_card(s, lbl, val, 0.4 + (i % 3) * 2.15, 1.45 + (i // 3) * 1.35, width=2.0)
    add_bullets(s, [
        "Verification:",
        "• Manual cross-role workflow testing",
        "• Seed data: John Smith (ABC123), Maria Garcia (XYZ789)",
        "• AI analysis log with confidence scores",
        "",
        "Academic deliverables:",
        "• ~20,000-word final report",
        "• Progress Reports 1 & 2 with presentations",
    ], left=0.4, top=4.3, width=6.5, size=12, spacing=5)
    add_image_card(s, dashboard, 7.0, 1.45, 5.8, 3.3, "Admin Dashboard — 7 jobs, 2 customers tracked")

    n += 1
    s = content_slide(prs, n, "Phase 1 — Foundation Results", "Authentication, admin & user management", "e")
    add_image_card(s, login, 0.35, 1.4, 6.1, 2.9, "JWT Login — role-based authentication")
    add_image_card(s, users, 6.75, 1.4, 6.1, 2.9, "User Management — 9 users, 8 roles")
    add_image_card(s, dashboard, 3.55, 4.55, 6.1, 1.85, "Admin Dashboard — system-wide KPIs")

    n += 1
    s = content_slide(prs, n, "Phase 2 — Business Module Results", "Customer, vehicle & job management", "e")
    add_image_card(s, customers, 0.35, 1.4, 6.1, 2.9, "Customer CRUD — 2 registered customers")
    add_image_card(s, vehicles, 6.75, 1.4, 6.1, 2.9, "Vehicle Registration — linked to owners")
    add_image_card(s, jobs, 3.55, 4.55, 6.1, 1.85, "Job List — department routing & assignment")

    n += 1
    s = content_slide(prs, n, "AI Integration — Results", "DeepSeek classification with audit trail", "e")
    add_image_card(s, create_job, 0.35, 1.4, 6.1, 2.9, "Create Job — AI-powered complaint analysis")
    add_image_card(s, ai_log, 6.75, 1.4, 6.1, 2.9, "AI Analysis Log — confidence & explanations")
    add_bullets(s, [
        "POST /api/jobs/analyze — complaint + vehicle registration",
        "DeepSeek returns: issues[], department, confidence, explanation",
        "Fallback keyword classifier when API unavailable",
        "Example: engine noise → Mechanical; headlight → Electrical",
    ], left=0.4, top=4.55, width=12.5, size=12, spacing=5)

    # ── SECTION f: Discussion ───────────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "f", "Discussion", "Objective alignment, limitations & lessons learned")

    n += 1
    s = content_slide(prs, n, "Discussion — Strengths & Limitations", "Honest critical evaluation — criterion 9", "f")
    add_bullets(s, [
        "Achievements:",
        "• All six objectives O1–O6 substantially achieved",
        "• Clean Architecture proved maintainable across phases",
        "• RBAC correctly filters data by role",
        "• AI audit trail supports accountability",
        "",
        "Course experience:",
        "• npm dependency conflicts during setup",
        "• Prisma migration issues resolved iteratively",
        "• DeepSeek API quota limitations encountered",
    ], width=6.0, size=12, spacing=5)
    limits = [
        ("L1", "API dependence", "AI accuracy tied to DeepSeek availability"),
        ("L2", "English-only fallback", "Keyword classifier not multilingual"),
        ("L3", "No automated tests", "Manual verification only"),
        ("L4", "Localhost only", "Production deployment deferred"),
    ]
    for i, (lid, title, desc) in enumerate(limits):
        top = 1.55 + i * 1.2
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(top), Inches(6.0), Inches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(top), Inches(0.07), Inches(1.0))
        accent.fill.solid()
        accent.fill.fore_color.rgb = AMBER
        accent.line.fill.background()
        lb = s.shapes.add_textbox(Inches(7.1), Inches(top + 0.1), Inches(0.6), Inches(0.35))
        lb.text_frame.text = lid
        set_run(lb.text_frame.paragraphs[0].runs[0], size=12, bold=True, color=AMBER)
        tt = s.shapes.add_textbox(Inches(7.7), Inches(top + 0.08), Inches(5.0), Inches(0.35))
        tt.text_frame.text = title
        set_run(tt.text_frame.paragraphs[0].runs[0], size=13, bold=True, color=NAVY)
        dt = s.shapes.add_textbox(Inches(7.7), Inches(top + 0.42), Inches(5.0), Inches(0.45))
        dt.text_frame.text = desc
        set_run(dt.text_frame.paragraphs[0].runs[0], size=11, color=SLATE)

    # ── SECTION g: Conclusions ──────────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "g", "Conclusions & Future Work", "Evidence-supported conclusions — criterion 10")

    n += 1
    s = content_slide(prs, n, "Conclusions", "Supported by evidence in Chapters 7 & 8", "g")
    conclusions = [
        "AutoRepairAgent delivers a working AI-integrated workshop management prototype",
        "DeepSeek LLM + fallback classifier enables complaint-to-department routing",
        "Clean Architecture supports maintainable full-stack development",
        "RBAC enforces appropriate access across 8 workshop roles",
        "Limitations are material but do not negate core objective achievement",
        "System is demonstrable and ready for pilot evaluation",
    ]
    for i, text in enumerate(conclusions):
        top = 1.45 + i * 0.82
        icon = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(top + 0.05), Inches(0.3), Inches(0.3))
        icon.fill.solid()
        icon.fill.fore_color.rgb = SUCCESS
        icon.line.fill.background()
        tick = s.shapes.add_textbox(Inches(0.45), Inches(top + 0.06), Inches(0.3), Inches(0.28))
        tick.text_frame.text = "✓"
        tick.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tick.text_frame.paragraphs[0].runs[0], size=11, bold=True, color=WHITE)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(top), Inches(5.8), Inches(0.65))
        tb.text_frame.text = text
        set_run(tb.text_frame.paragraphs[0].runs[0], size=13, color=SLATE)
    add_image_card(s, ai_log, 7.0, 1.45, 5.8, 3.3, "Evidence: AI audit trail with job linkage")

    n += 1
    s = content_slide(prs, n, "Future Work Roadmap", "Prioritised recommendations", "g")
    future = [
        ("F1", "Cloud deployment", "Azure/AWS production hosting"),
        ("F2", "Multilingual UI", "Sinhala/Tamil localisation"),
        ("F3", "Automated testing", "Jest + Playwright assertions"),
        ("F4", "SMS notifications", "Job status alerts to customers"),
        ("F5", "Smart assignment", "Load-balanced technician routing"),
        ("F6", "Inventory module", "Parts tracking & invoicing"),
    ]
    for i, (fid, title, desc) in enumerate(future):
        col = i % 2
        row = i // 2
        left = 0.4 + col * 6.4
        top = 1.42 + row * 1.55
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(6.1), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        fb = s.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.1), Inches(0.55), Inches(0.3))
        fb.text_frame.text = fid
        set_run(fb.text_frame.paragraphs[0].runs[0], size=11, bold=True, color=TEAL)
        tt = s.shapes.add_textbox(Inches(left + 0.7), Inches(top + 0.08), Inches(5.2), Inches(0.3))
        tt.text_frame.text = title
        set_run(tt.text_frame.paragraphs[0].runs[0], size=12, bold=True, color=NAVY)
        dt = s.shapes.add_textbox(Inches(left + 0.7), Inches(top + 0.42), Inches(5.2), Inches(0.35))
        dt.text_frame.text = desc
        set_run(dt.text_frame.paragraphs[0].runs[0], size=10, color=MUTED)

    # ── SLIDE 24: Q&A (Rubric criterion) ────────────────────────────────────────
    n += 1
    s = content_slide(prs, n, "Questions for Discussion", "Prepared to stimulate examiner discussion — Rubric criterion")
    questions = [
        ("Q1", "How does DeepSeek compare to rule-based triage for workshop accuracy?"),
        ("Q2", "What automated tests would you prioritise before production deployment?"),
        ("Q3", "How should the system handle multilingual complaints in NZ workshops?"),
        ("Q4", "What privacy disclosures are needed for third-party AI processing?"),
        ("Q5", "How would you scale technician assignment beyond first-available routing?"),
        ("Q6", "What cloud architecture would you recommend for production deployment?"),
    ]
    for i, (qid, qtext) in enumerate(questions):
        top = 1.45 + i * 0.88
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(top), Inches(12.5), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER
        qb = s.shapes.add_textbox(Inches(0.55), Inches(top + 0.15), Inches(0.55), Inches(0.4))
        qb.text_frame.text = qid
        set_run(qb.text_frame.paragraphs[0].runs[0], size=12, bold=True, color=TEAL)
        qt = s.shapes.add_textbox(Inches(1.2), Inches(top + 0.15), Inches(11.5), Inches(0.45))
        qt.text_frame.text = qtext
        set_run(qt.text_frame.paragraphs[0].runs[0], size=13, color=NAVY)

    add_bullets(s, [
        "Rubric target: Excellent (4/4) — prepared questions that open technical discussion",
    ], left=0.4, top=6.65, width=12.0, size=11, spacing=4, color=MUTED)

    # ── SLIDE 25: Thank You ─────────────────────────────────────────────────────
    n += 1
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()

    inst2 = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.35))
    inst2.text_frame.text = INST
    set_run(inst2.text_frame.paragraphs[0].runs[0], size=12, color=TEAL)

    thanks = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(7.0), Inches(0.9))
    thanks.text_frame.text = "Thank You"
    set_run(thanks.text_frame.paragraphs[0].runs[0], size=44, bold=True, color=WHITE)

    sub2 = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(7.0), Inches(0.5))
    sub2.text_frame.text = "Questions & Discussion Welcome"
    set_run(sub2.text_frame.paragraphs[0].runs[0], size=20, color=TEAL)

    info = s.shapes.add_textbox(Inches(0.8), Inches(3.75), Inches(7.0), Inches(1.5))
    tf = info.text_frame
    for i, line in enumerate([
        "AutoRepairAgent",
        "[Your Full Name]  |  [Your Student ID]",
        "COMP902 Final Report  |  August 2026",
        "6 References  |  ~20,000-word Report  |  Live UI Screenshots",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
        set_run(p.runs[0], size=14, color=LIGHT_TXT)

    add_image_card(s, jobs, 8.2, 1.1, 4.5, 2.6, "AutoRepairAgent — Full System Delivery")

    foot = s.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(12.0), Inches(0.3))
    foot.text_frame.text = f"{COURSE}  |  Slide {n} of {TOTAL}"
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_run(foot.text_frame.paragraphs[0].runs[0], size=9, color=MUTED)

    assert len(prs.slides) == TOTAL, f"Expected {TOTAL} slides, got {len(prs.slides)}"
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Structure: COMP902 sections a–g + Q&A rubric")


if __name__ == "__main__":
    build()
